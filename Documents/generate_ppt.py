"""
Generate ValiData Overview Presentation (2 slides)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──
DARK_NAVY   = RGBColor(0x0F, 0x1A, 0x2E)
MID_NAVY    = RGBColor(0x1E, 0x3A, 0x5F)
ACCENT_BLUE = RGBColor(0x00, 0x9D, 0xE0)
LIGHT_BLUE  = RGBColor(0xDC, 0xE9, 0xF8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE      = RGBColor(0xF5, 0x7C, 0x00)
GREEN       = RGBColor(0x2E, 0xCC, 0x71)
RED_SOFT    = RGBColor(0xE7, 0x4C, 0x3C)
GOLD        = RGBColor(0xF1, 0xC4, 0x0F)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════════════
# Helper functions
# ═══════════════════════════════════════════════════════════════════

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=WHITE, spacing=Pt(8), font_name="Calibri", icon_color=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
    return txBox

def add_icon_bullet_card(slide, left, top, width, items, card_fill, icon_col, text_col=WHITE, font_size=12, card_height=None):
    """Add a card with icon-prefixed bullet items"""
    item_h = Pt(font_size) * 2.8
    h = card_height or int(len(items) * item_h + Inches(0.4))
    card = add_rect(slide, left, top, width, h, card_fill)
    
    y_offset = top + Inches(0.15)
    for icon, text in items:
        line = f"{icon}  {text}"
        add_text_box(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.35), line, font_size=font_size, color=text_col)
        y_offset += Inches(0.38)
    return card


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Problem Statement + Our Solution
# ═══════════════════════════════════════════════════════════════════

slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, DARK_NAVY)

# ── Top banner ──
add_rect(slide1, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MID_NAVY)
add_text_box(slide1, Inches(0.6), Inches(0.2), Inches(8), Inches(0.55), "ValiData", font_size=32, bold=True, color=WHITE)
add_text_box(slide1, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35), "Data Quality & Observability Control Plane for Snowflake & Databricks", font_size=15, color=LIGHT_GRAY)

# ── Accent line ──
line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# ═══════ LEFT HALF — Problem Statement ═══════
add_text_box(slide1, Inches(0.6), Inches(1.4), Inches(5.5), Inches(0.45), "⚠  THE PROBLEM", font_size=20, bold=True, color=ORANGE)

problem_items = [
    ("🔓", "Data moves outside the warehouse for quality checks — creating security risks & compliance gaps"),
    ("💸", "High egress costs when extracting data to external DQ tools for validation"),
    ("🔍", "No unified view — teams juggle multiple tools for profiling, lineage, anomaly detection & governance"),
    ("⏱", "Manual, ad-hoc quality checks — no scheduling, no automation, no audit trail"),
    ("🧩", "Siloed metadata — data catalog, quality scores, and lineage live in separate systems"),
]

y = Inches(1.95)
for icon, text in problem_items:
    # Icon circle
    circ = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.02), Inches(0.35), Inches(0.35))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x2A, 0x1A, 0x1A)
    circ.line.color.rgb = RED_SOFT
    circ.line.width = Pt(1.5)
    add_text_box(slide1, Inches(0.72), y + Inches(0.03), Inches(0.35), Inches(0.35), icon, font_size=14, alignment=PP_ALIGN.CENTER)
    
    # Text
    add_text_box(slide1, Inches(1.2), y, Inches(5.2), Inches(0.45), text, font_size=12, color=LIGHT_GRAY)
    y += Inches(0.55)

# Bottom summary box for problem
add_rect(slide1, Inches(0.6), y + Inches(0.15), Inches(5.8), Inches(0.5), RGBColor(0x3A, 0x15, 0x15), border_color=RED_SOFT, border_width=Pt(1))
add_text_box(slide1, Inches(0.8), y + Inches(0.2), Inches(5.5), Inches(0.4),
    "Result: Poor data trust, delayed decisions, wasted compute spend",
    font_size=12, bold=True, color=RED_SOFT, alignment=PP_ALIGN.CENTER)

# ═══════ Vertical Divider ═══════
div = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.3), Pt(2), Inches(5.8))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x30, 0x50, 0x70)
div.line.fill.background()

# ═══════ RIGHT HALF — Our Solution ═══════
add_text_box(slide1, Inches(6.95), Inches(1.4), Inches(5.5), Inches(0.45), "✅  OUR SOLUTION", font_size=20, bold=True, color=GREEN)

# Solution sub-header
add_text_box(slide1, Inches(6.95), Inches(1.85), Inches(5.8), Inches(0.45),
    "\"Pushdown Architecture\" — SQL runs inside the warehouse, data never leaves",
    font_size=12, bold=True, color=ACCENT_BLUE)

# Solution cards
card_w = Inches(5.8)

# Card 1: DQ Engine
add_rect(slide1, Inches(6.95), Inches(2.35), card_w, Inches(0.9), RGBColor(0x15, 0x2A, 0x40), border_color=ACCENT_BLUE, border_width=Pt(1))
add_text_box(slide1, Inches(7.15), Inches(2.4), card_w, Inches(0.3), "📊  Data Quality Engine", font_size=13, bold=True, color=ACCENT_BLUE)
add_text_box(slide1, Inches(7.15), Inches(2.7), Inches(5.5), Inches(0.5),
    "Define rules (Null, Unique, Range, Pattern, Blank) → ValiData generates SQL → pushes to Snowflake/Databricks → returns scores & anomalies",
    font_size=11, color=LIGHT_GRAY)

# Card 2: AI-Powered
add_rect(slide1, Inches(6.95), Inches(3.35), card_w, Inches(0.9), RGBColor(0x15, 0x2A, 0x40), border_color=GOLD, border_width=Pt(1))
add_text_box(slide1, Inches(7.15), Inches(3.4), card_w, Inches(0.3), "🤖  AI-Powered Intelligence (No External API Keys)", font_size=13, bold=True, color=GOLD)
add_text_box(slide1, Inches(7.15), Inches(3.7), Inches(5.5), Inches(0.5),
    "Snowflake Cortex (Mistral-large) & Databricks AI Functions (Llama 3) — LLMs run inside the warehouse for rule suggestions, chat agent, table summaries",
    font_size=11, color=LIGHT_GRAY)

# Card 3: Observability
add_rect(slide1, Inches(6.95), Inches(4.35), card_w, Inches(0.9), RGBColor(0x15, 0x2A, 0x40), border_color=GREEN, border_width=Pt(1))
add_text_box(slide1, Inches(7.15), Inches(4.4), card_w, Inches(0.3), "🔗  Full Observability Suite", font_size=13, bold=True, color=GREEN)
add_text_box(slide1, Inches(7.15), Inches(4.7), Inches(5.5), Inches(0.5),
    "Live Data Catalog • Column Profiling • Automated Lineage Inference • Usage Analytics • Anomaly Detection • Scheduled DQ Runs",
    font_size=11, color=LIGHT_GRAY)

# Tech Stack bar
add_rect(slide1, Inches(6.95), Inches(5.45), card_w, Inches(1.4), RGBColor(0x10, 0x20, 0x35), border_color=RGBColor(0x30, 0x50, 0x70), border_width=Pt(1))
add_text_box(slide1, Inches(7.15), Inches(5.5), card_w, Inches(0.3), "🛠  Technology Stack", font_size=13, bold=True, color=WHITE)

tech_items = [
    "Frontend:     React 19  •  TypeScript  •  Vite 8  •  React Flow",
    "Backend:      FastAPI (Python)  •  Uvicorn",
    "Connectors:  snowflake-connector-python  •  databricks-sql-connector",
    "AI / LLM:      Snowflake Cortex (Mistral)  •  Databricks AI (Llama 3)",
    "Database:     PostgreSQL (Neon)  /  SQLite  (auto-detects)",
]
ty = Inches(5.82)
for t in tech_items:
    add_text_box(slide1, Inches(7.25), ty, Inches(5.4), Inches(0.22), t, font_size=10, color=LIGHT_GRAY)
    ty += Inches(0.22)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: Benefits
# ═══════════════════════════════════════════════════════════════════

slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, DARK_NAVY)

# ── Top banner ──
add_rect(slide2, Inches(0), Inches(0), Inches(13.333), Inches(1.1), MID_NAVY)
add_text_box(slide2, Inches(0.6), Inches(0.2), Inches(8), Inches(0.55), "ValiData — Key Benefits", font_size=32, bold=True, color=WHITE)
add_text_box(slide2, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35), "Why ValiData is the right choice for enterprise data quality", font_size=15, color=LIGHT_GRAY)

line2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Pt(3))
line2.fill.solid()
line2.fill.fore_color.rgb = ACCENT_BLUE
line2.line.fill.background()

# ── Benefit cards (2x3 grid) ──
benefits = [
    {
        "icon": "🔒",
        "title": "Zero Data Movement",
        "desc": "Raw data NEVER leaves Snowflake / Databricks. SQL is pushed into the warehouse. Only metadata, counts, and anomaly flags are returned.",
        "color": ACCENT_BLUE,
    },
    {
        "icon": "💰",
        "title": "Zero Egress Cost",
        "desc": "No data extraction means no egress charges. Uses existing warehouse compute — no separate DQ infrastructure to provision or pay for.",
        "color": GREEN,
    },
    {
        "icon": "🤖",
        "title": "Native AI (No API Keys)",
        "desc": "Leverages Snowflake Cortex & Databricks AI Functions. LLMs run inside the warehouse for rule suggestions, chat, and table summaries.",
        "color": GOLD,
    },
    {
        "icon": "📊",
        "title": "Unified Observability",
        "desc": "One platform for Data Catalog, Column Profiling, DQ Rules, Lineage, Usage Analytics, and Anomaly Detection — no tool sprawl.",
        "color": ORANGE,
    },
    {
        "icon": "⏱",
        "title": "Automated & Scheduled",
        "desc": "Schedule recurring DQ runs (every 5 min to monthly). Automatic anomaly detection with audit trail. No manual intervention needed.",
        "color": RGBColor(0xAF, 0x7A, 0xC5),
    },
    {
        "icon": "🚀",
        "title": "Cloud-Native & Fast",
        "desc": "Frontend on Vercel (global CDN), Backend on Render (auto-scaling), DB on Neon (serverless PostgreSQL). Deploy in minutes, not weeks.",
        "color": RGBColor(0x00, 0xD2, 0xFF),
    },
]

card_w = Inches(3.9)
card_h = Inches(2.4)
x_positions = [Inches(0.6), Inches(4.75), Inches(8.9)]
y_positions = [Inches(1.45), Inches(4.15)]

for idx, b in enumerate(benefits):
    col = idx % 3
    row = idx // 3
    x = x_positions[col]
    y = y_positions[row]
    
    # Card background
    card = add_rect(slide2, x, y, card_w, card_h, RGBColor(0x15, 0x2A, 0x40), border_color=b["color"], border_width=Pt(1.5))
    
    # Icon circle
    circ = slide2.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circ.fill.solid()
    circ.fill.fore_color.rgb = RGBColor(0x0A, 0x15, 0x25)
    circ.line.color.rgb = b["color"]
    circ.line.width = Pt(2)
    add_text_box(slide2, x + Inches(0.2), y + Inches(0.22), Inches(0.5), Inches(0.5),
        b["icon"], font_size=20, alignment=PP_ALIGN.CENTER)
    
    # Title
    add_text_box(slide2, x + Inches(0.85), y + Inches(0.25), Inches(2.8), Inches(0.4),
        b["title"], font_size=16, bold=True, color=b["color"])
    
    # Description
    add_text_box(slide2, x + Inches(0.25), y + Inches(0.85), Inches(3.4), Inches(1.4),
        b["desc"], font_size=12, color=LIGHT_GRAY)

# ── Bottom tagline ──
add_rect(slide2, Inches(2.5), Inches(6.8), Inches(8.333), Inches(0.5), RGBColor(0x10, 0x20, 0x35), border_color=ACCENT_BLUE, border_width=Pt(1))
add_text_box(slide2, Inches(2.5), Inches(6.83), Inches(8.333), Inches(0.45),
    "ValiData  —  Enterprise Trust, Native Performance, Zero Data Movement",
    font_size=14, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = r"c:\Users\Damini\Documents\BristleCone_POC\ValiData\Documents\ValiData_Overview.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
